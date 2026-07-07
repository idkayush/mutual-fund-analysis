from pathlib import Path
import pandas as pd
from reportlab.lib.pagesizes import A4
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle, PageBreak
from pptx import Presentation
from pptx.util import Inches, Pt


BASE_DIR = Path(".")
REPORTS_DIR = BASE_DIR / "reports"
CHARTS_DIR = REPORTS_DIR / "charts"
SCREENSHOT_DIR = REPORTS_DIR / "dashboard_screenshots"

FINAL_REPORT = REPORTS_DIR / "Final_Report.pdf"
PRESENTATION = REPORTS_DIR / "Presentation.pptx"


def safe_read_csv(path):
    if path.exists():
        return pd.read_csv(path)
    return pd.DataFrame()


def get_top_funds():
    scorecard = safe_read_csv(REPORTS_DIR / "fund_scorecard.csv")
    if scorecard.empty:
        return []
    cols = ["scheme_name", "fund_house", "category", "sharpe_ratio", "fund_score"]
    cols = [c for c in cols if c in scorecard.columns]
    return scorecard.sort_values("fund_score", ascending=False)[cols].head(5).values.tolist()


def build_pdf():
    styles = getSampleStyleSheet()
    story = []

    doc = SimpleDocTemplate(
        str(FINAL_REPORT),
        pagesize=A4,
        rightMargin=36,
        leftMargin=36,
        topMargin=36,
        bottomMargin=36
    )

    title = styles["Title"]
    heading = styles["Heading1"]
    subheading = styles["Heading2"]
    body = styles["BodyText"]

    story.append(Paragraph("Bluestock Mutual Fund Analytics Capstone", title))
    story.append(Spacer(1, 12))
    story.append(Paragraph("Final Project Report", subheading))
    story.append(Spacer(1, 18))

    story.append(Paragraph("1. Executive Summary", heading))
    story.append(Paragraph(
        "This project presents a complete mutual fund analytics pipeline covering ETL, data cleaning, SQLite database design, exploratory data analysis, performance analytics, advanced risk analytics, investor cohort analysis, recommender logic, and an interactive Power BI dashboard.",
        body
    ))
    story.append(Spacer(1, 12))

    story.append(Paragraph("2. Project Objectives", heading))
    objectives = [
        "Clean and validate mutual fund NAV, AUM, SIP, transaction, performance, benchmark, and portfolio datasets.",
        "Design a SQLite star schema for structured analytics.",
        "Perform EDA using 15+ charts and document important business insights.",
        "Compute CAGR, Sharpe Ratio, Sortino Ratio, Alpha, Beta, Tracking Error, Maximum Drawdown, VaR, and CVaR.",
        "Build an interactive Power BI dashboard with four analytical pages.",
        "Develop advanced analytics including investor cohorts, SIP continuity analysis, sector concentration, and fund recommender."
    ]
    for obj in objectives:
        story.append(Paragraph(f"• {obj}", body))
    story.append(Spacer(1, 12))

    story.append(Paragraph("3. Data Pipeline and Cleaning", heading))
    story.append(Paragraph(
        "The ETL pipeline loads raw CSV datasets, validates schema consistency, parses date columns, removes duplicates, validates positive NAV and transaction amounts, standardizes transaction types, checks KYC status values, and creates cleaned CSVs in the processed data folder.",
        body
    ))
    story.append(Spacer(1, 12))

    story.append(Paragraph("4. SQLite Database Design", heading))
    story.append(Paragraph(
        "A star schema was designed using dimension tables such as dim_fund and dim_date, and fact tables such as fact_nav, fact_transactions, fact_performance, and fact_aum. SQL schema and analytical queries are stored in the sql folder.",
        body
    ))
    story.append(Spacer(1, 12))

    story.append(Paragraph("5. EDA Highlights", heading))
    eda_points = [
        "NAV trends show broad growth across schemes with visible market movements during 2023 and 2024.",
        "SIP inflows show strong growth from 2022 to 2025.",
        "AUM analysis highlights dominance of major AMCs.",
        "Investor analytics shows transaction concentration across states, age groups, and city tiers.",
        "Sector allocation and category inflow charts reveal portfolio and market preference patterns."
    ]
    for point in eda_points:
        story.append(Paragraph(f"• {point}", body))

    chart_paths = [
        CHARTS_DIR / "03_aum_growth_by_fund_house.png",
        CHARTS_DIR / "04_sip_inflow_time_series.png",
        CHARTS_DIR / "09_sip_amount_by_state.png",
        CHARTS_DIR / "12_nav_return_correlation_matrix.png",
        CHARTS_DIR / "benchmark_comparison_top5_vs_nifty.png"
    ]

    for chart in chart_paths:
        if chart.exists():
            story.append(Spacer(1, 12))
            story.append(Image(str(chart), width=480, height=270))

    story.append(PageBreak())

    story.append(Paragraph("6. Performance Analytics", heading))
    story.append(Paragraph(
        "Performance analytics included CAGR computation, Sharpe Ratio, Sortino Ratio, Alpha, Beta, Tracking Error, and Maximum Drawdown. A composite fund scorecard was created using weighted ranks across return, risk-adjusted return, alpha, expense ratio, and drawdown.",
        body
    ))
    story.append(Spacer(1, 12))

    top_funds = get_top_funds()
    if top_funds:
        story.append(Paragraph("Top 5 Funds by Composite Score", subheading))
        table_data = [["Scheme", "Fund House", "Category", "Sharpe", "Score"]]
        for row in top_funds:
            clean_row = []
            for value in row:
                if isinstance(value, float):
                    clean_row.append(round(value, 3))
                else:
                    clean_row.append(str(value)[:30])
            table_data.append(clean_row)

        table = Table(table_data, repeatRows=1)
        table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0B5FFF")),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, -1), 8),
            ("VALIGN", (0, 0), (-1, -1), "TOP"),
        ]))
        story.append(table)

    story.append(Spacer(1, 12))
    story.append(Paragraph("7. Advanced Analytics", heading))
    advanced_points = [
        "Historical VaR and CVaR were computed at 95 percent confidence for all schemes.",
        "Rolling 90-day Sharpe Ratio was plotted for key funds.",
        "Investor cohorts were grouped by first transaction year.",
        "SIP continuity analysis identified investors with gaps greater than 35 days.",
        "A simple recommender ranks funds by Sharpe Ratio within the selected risk appetite.",
        "Sector HHI concentration was calculated to identify concentrated equity portfolios."
    ]
    for point in advanced_points:
        story.append(Paragraph(f"• {point}", body))

    story.append(Spacer(1, 12))
    story.append(Paragraph("8. Power BI Dashboard", heading))
    story.append(Paragraph(
        "The Power BI dashboard contains four pages: Industry Overview, Fund Performance, Investor Analytics, and SIP & Market Trends. It includes slicers, KPI cards, line charts, bar charts, donut charts, heatmap matrix, scorecard table, and benchmark comparison visuals.",
        body
    ))

    screenshots = [
        SCREENSHOT_DIR / "page1_industry_overview.png",
        SCREENSHOT_DIR / "page2_fund_performance.png",
        SCREENSHOT_DIR / "page3_investor_analytics.png",
        SCREENSHOT_DIR / "page4_sip_market_trends.png"
    ]

    for ss in screenshots:
        if ss.exists():
            story.append(Spacer(1, 12))
            story.append(Image(str(ss), width=480, height=270))

    story.append(PageBreak())

    story.append(Paragraph("9. Final Findings", heading))
    findings = [
        "The mutual fund industry shows strong growth in SIP participation and folio count.",
        "Fund performance varies significantly by category, volatility, and risk-adjusted return.",
        "Composite scoring provides a more balanced ranking than return-only comparison.",
        "Investor behaviour differs across geography, city tiers, and age groups.",
        "Advanced metrics such as VaR, CVaR, drawdown, and HHI improve risk visibility."
    ]
    for finding in findings:
        story.append(Paragraph(f"• {finding}", body))

    story.append(Spacer(1, 12))
    story.append(Paragraph("10. Conclusion", heading))
    story.append(Paragraph(
        "This capstone demonstrates an end-to-end analytics workflow from raw mutual fund data to business-ready insights, risk analytics, fund ranking, recommender logic, and an interactive dashboard. The project is structured for reproducibility and professional submission.",
        body
    ))

    doc.build(story)
    print(f"Saved: {FINAL_REPORT}")


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for bullet in bullets:
        p = body.add_paragraph()
        p.text = bullet
        p.font.size = Pt(20)


def add_image_slide(prs, title, image_path):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(0.7), Inches(1.3), width=Inches(8.8))


def build_ppt():
    prs = Presentation()

    add_title_slide(
        prs,
        "Bluestock Mutual Fund Analytics",
        "ETL, EDA, Performance Analytics, Risk Analytics, Recommender and Power BI Dashboard"
    )

    add_bullet_slide(prs, "Project Scope", [
        "Built end-to-end mutual fund analytics pipeline",
        "Cleaned and validated NAV, AUM, SIP, investor, performance, benchmark and portfolio data",
        "Designed SQLite star schema and analytical SQL queries",
        "Created EDA, performance analytics, advanced analytics and Power BI dashboard"
    ])

    add_bullet_slide(prs, "Technical Architecture", [
        "Raw CSV and live NAV data ingestion",
        "Python-based cleaning and feature engineering",
        "SQLite database with fact and dimension tables",
        "Jupyter notebooks for EDA and analytics",
        "Power BI dashboard for interactive insights"
    ])

    add_bullet_slide(prs, "Key Metrics Computed", [
        "CAGR for 1-year, 3-year and 5-year periods",
        "Sharpe Ratio and Sortino Ratio",
        "Alpha and Beta using benchmark regression",
        "Tracking Error and Maximum Drawdown",
        "VaR, CVaR, rolling Sharpe and sector HHI"
    ])

    add_image_slide(prs, "AUM Growth by Fund House", CHARTS_DIR / "03_aum_growth_by_fund_house.png")
    add_image_slide(prs, "SIP Inflow Trend", CHARTS_DIR / "04_sip_inflow_time_series.png")
    add_image_slide(prs, "Benchmark Comparison", CHARTS_DIR / "benchmark_comparison_top5_vs_nifty.png")
    add_image_slide(prs, "Dashboard: Industry Overview", SCREENSHOT_DIR / "page1_industry_overview.png")
    add_image_slide(prs, "Dashboard: Fund Performance", SCREENSHOT_DIR / "page2_fund_performance.png")
    add_image_slide(prs, "Dashboard: Investor Analytics", SCREENSHOT_DIR / "page3_investor_analytics.png")
    add_image_slide(prs, "Dashboard: SIP & Market Trends", SCREENSHOT_DIR / "page4_sip_market_trends.png")

    add_bullet_slide(prs, "Advanced Analytics Insights", [
        "VaR and CVaR identify funds with higher downside risk",
        "Investor cohort analysis shows year-wise investment behaviour",
        "SIP continuity analysis flags at-risk investors",
        "Sector HHI identifies concentrated equity fund portfolios",
        "Risk appetite recommender suggests top funds by Sharpe Ratio"
    ])

    add_bullet_slide(prs, "Conclusion", [
        "End-to-end analytics solution completed",
        "Dashboard enables interactive business exploration",
        "Risk and performance metrics improve fund comparison",
        "Project structure is reproducible and submission-ready"
    ])

    prs.save(PRESENTATION)
    print(f"Saved: {PRESENTATION}")


if __name__ == "__main__":
    REPORTS_DIR.mkdir(exist_ok=True)
    build_pdf()
    build_ppt()
    print("Final report and presentation generated successfully.")
